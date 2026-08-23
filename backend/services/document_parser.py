import io
import re
from typing import List, Dict, Any
from pypdf import PdfReader
import docx
from pptx import Presentation

class DocumentParser:
    @staticmethod
    def extract_text(file_bytes: bytes, filename: str, file_type: str) -> Dict[str, Any]:
        """
        Extract structured text sections with metadata from various file formats:
        - PDF (.pdf)
        - DOCX (.docx, .doc)
        - PPTX (.pptx, .ppt)
        - TXT (.txt)
        - Markdown (.md)
        """
        filename_lower = filename.lower()
        extracted_pages: List[Dict[str, Any]] = []
        raw_text_parts: List[str] = []

        try:
            if filename_lower.endswith(".pdf"):
                reader = PdfReader(io.BytesIO(file_bytes))
                for idx, page in enumerate(reader.pages):
                    text = page.extract_text() or ""
                    cleaned = DocumentParser.clean_text(text)
                    if cleaned:
                        extracted_pages.append({
                            "page_number": idx + 1,
                            "section_title": f"Page {idx + 1}",
                            "content": cleaned
                        })
                        raw_text_parts.append(cleaned)

            elif filename_lower.endswith((".docx", ".doc")):
                doc = docx.Document(io.BytesIO(file_bytes))
                current_section = "Introduction / Overview"
                section_paragraphs = []
                
                for p in doc.paragraphs:
                    p_text = p.text.strip()
                    if not p_text:
                        continue
                    # Check if paragraph is heading
                    if p.style and hasattr(p.style, "name") and p.style.name and p.style.name.startswith("Heading"):
                        if section_paragraphs:
                            content_str = "\n".join(section_paragraphs)
                            extracted_pages.append({
                                "page_number": 1,
                                "section_title": current_section,
                                "content": content_str
                            })
                            raw_text_parts.append(content_str)
                            section_paragraphs = []
                        current_section = p_text
                    else:
                        section_paragraphs.append(p_text)

                if section_paragraphs:
                    content_str = "\n".join(section_paragraphs)
                    extracted_pages.append({
                        "page_number": 1,
                        "section_title": current_section,
                        "content": content_str
                    })
                    raw_text_parts.append(content_str)

            elif filename_lower.endswith((".pptx", ".ppt")):
                prs = Presentation(io.BytesIO(file_bytes))
                for idx, slide in enumerate(prs.slides):
                    slide_texts = []
                    slide_title = f"Slide {idx + 1}"
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                t = paragraph.text.strip()
                                if t:
                                    slide_texts.append(t)
                    if slide_texts:
                        # Use first line as title if possible
                        if len(slide_texts) > 0 and len(slide_texts[0]) < 80:
                            slide_title = slide_texts[0]
                        joined = "\n".join(slide_texts)
                        cleaned = DocumentParser.clean_text(joined)
                        if cleaned:
                            extracted_pages.append({
                                "page_number": idx + 1,
                                "section_title": slide_title,
                                "content": cleaned
                            })
                            raw_text_parts.append(cleaned)

            elif filename_lower.endswith((".txt", ".md")):
                text = file_bytes.decode("utf-8", errors="replace")
                cleaned = DocumentParser.clean_text(text)
                
                # Split by markdown headers if present
                sections = re.split(r'(?m)^#{1,3}\s+(.+)$', cleaned)
                if len(sections) > 1:
                    # sections[0] is preamble, then header1, content1, header2, content2...
                    if sections[0].strip():
                        extracted_pages.append({
                            "page_number": 1,
                            "section_title": "Overview",
                            "content": sections[0].strip()
                        })
                    for i in range(1, len(sections), 2):
                        sec_title = sections[i].strip()
                        sec_body = sections[i+1].strip() if i+1 < len(sections) else ""
                        if sec_body:
                            extracted_pages.append({
                                "page_number": (i // 2) + 1,
                                "section_title": sec_title,
                                "content": f"# {sec_title}\n{sec_body}"
                            })
                            raw_text_parts.append(sec_body)
                else:
                    extracted_pages.append({
                        "page_number": 1,
                        "section_title": "Document Content",
                        "content": cleaned
                    })
                    raw_text_parts.append(cleaned)

            else:
                # Generic text fallback
                text = file_bytes.decode("utf-8", errors="replace")
                cleaned = DocumentParser.clean_text(text)
                extracted_pages.append({
                    "page_number": 1,
                    "section_title": "General Text",
                    "content": cleaned
                })
                raw_text_parts.append(cleaned)

        except Exception as e:
            # Fallback error handling
            print(f"Error parsing document {filename}: {e}")
            text = file_bytes.decode("utf-8", errors="ignore")
            cleaned = DocumentParser.clean_text(text)
            if cleaned:
                extracted_pages.append({
                    "page_number": 1,
                    "section_title": "Raw Document Content",
                    "content": cleaned
                })
                raw_text_parts.append(cleaned)

        full_text = "\n\n".join(raw_text_parts)
        is_empty = len(full_text.strip()) == 0

        return {
            "is_empty": is_empty,
            "pages": extracted_pages,
            "full_text": full_text,
            "total_chars": len(full_text)
        }

    @staticmethod
    def clean_text(text: str) -> str:
        """Removes duplicate whitespace, non-printable characters and trims lines."""
        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
        text = re.sub(r'\r\n', '\n', text)
        text = re.sub(r'[ \t]+', ' ', text)
        text = re.sub(r'\n{3,}', '\n\n', text)
        return text.strip()

    @staticmethod
    def chunk_document(pages: List[Dict[str, Any]], chunk_size: int = 800, chunk_overlap: int = 150) -> List[Dict[str, Any]]:
        """
        Creates semantic chunks preserving section metadata and page numbers.
        """
        chunks: List[Dict[str, Any]] = []
        chunk_idx = 0

        for page_data in pages:
            page_num = page_data.get("page_number", 1)
            section_title = page_data.get("section_title", "General")
            content = page_data.get("content", "")

            if not content:
                continue

            # If page content is short, make it a single chunk
            if len(content) <= chunk_size:
                chunks.append({
                    "chunk_index": chunk_idx,
                    "page_number": page_num,
                    "section_title": section_title,
                    "content": content
                })
                chunk_idx += 1
                continue

            # Split content by paragraphs or sentences
            paragraphs = content.split("\n")
            current_chunk_text = ""

            for p in paragraphs:
                p_clean = p.strip()
                if not p_clean:
                    continue

                if len(current_chunk_text) + len(p_clean) < chunk_size:
                    current_chunk_text += ("\n" if current_chunk_text else "") + p_clean
                else:
                    if current_chunk_text:
                        chunks.append({
                            "chunk_index": chunk_idx,
                            "page_number": page_num,
                            "section_title": section_title,
                            "content": current_chunk_text
                        })
                        chunk_idx += 1
                        # Retain overlap from end of current chunk
                        current_chunk_text = current_chunk_text[-chunk_overlap:] + "\n" + p_clean
                    else:
                        # Single large paragraph - slice by chunk_size
                        for i in range(0, len(p_clean), chunk_size - chunk_overlap):
                            slice_text = p_clean[i:i + chunk_size]
                            chunks.append({
                                "chunk_index": chunk_idx,
                                "page_number": page_num,
                                "section_title": section_title,
                                "content": slice_text
                            })
                            chunk_idx += 1
                        current_chunk_text = ""

            if current_chunk_text.strip():
                chunks.append({
                    "chunk_index": chunk_idx,
                    "page_number": page_num,
                    "section_title": section_title,
                    "content": current_chunk_text.strip()
                })
                chunk_idx += 1

        return chunks
